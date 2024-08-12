'''
1.- Leer plantilla base
2.- Reemplazar texto
3.- Reemplazar imagenes
4.- Guardar archivo
'''

# %% Librerias
# ----------------------------------------------------------------------------------------------------------------------
from pptx import Presentation

file_path = 'ppt_test.pptx'
wildcard_replacement = {
    "%title%": "Monitoreo ventas trismetral",
    "%nombre%": "Richard Ramos",
    "%subtitulo1%": "Resultados Ventas 2018 - 2021"
}

rectangle_replacements = {
    "IMG1": "image1.png"
}
# Funcion principal
def process_ppt(file, wildcard_text, rectangle_img):
    ppt = Presentation(file)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text in wildcard_text:
                replace_text(shape, wildcard_text)
            elif shape.shape_type == 1:
                replace_rect_img(slide, shape, rectangle_img)

    ppt.save('salida.pptx')

# Funciones de ayuda
def replace_text(shape, wildcard):
    shape.text = wildcard[shape.text]

def replace_rect_img(slide, shape, rectangle):
    if shape.text in rectangle:
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height

        slide.shapes._spTree.remove(shape._element)
        slide.shapes.add_picture(rectangle[shape.text], left, top, width, height)


process_ppt(file_path, wildcard_replacement, rectangle_replacements)





#%%
